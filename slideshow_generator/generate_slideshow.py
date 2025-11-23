"""
Script untuk Generate Slideshow PowerPoint Yudisium
Oleh: Claude AI Assistant
Tanggal: 22 November 2025

Script ini akan:
1. Download foto mahasiswa dari Google Drive
2. Generate PowerPoint dengan 2 mahasiswa per slide
"""

import pandas as pd
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import os
import sys
from time import sleep
from io import BytesIO

# ============================================================================
# KONFIGURASI
# ============================================================================

EXCEL_FILE = 'Formulir_Lapor_Diri_Mahasiswa_PPG_2025_Batch_3.xlsx'
FOTO_FOLDER = 'foto_mahasiswa'
OUTPUT_PPTX = 'Slideshow_Yudisium_2025.pptx'

# Ukuran slide (16:9)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

# Posisi dan ukuran untuk 2 mahasiswa per slide
# Mahasiswa 1 (kiri)
FOTO1_LEFT = Inches(0.5)
FOTO1_TOP = Inches(0.5)
FOTO1_WIDTH = Inches(3.5)
FOTO1_HEIGHT = Inches(3.5)

# Mahasiswa 2 (kanan)
FOTO2_LEFT = Inches(6.0)
FOTO2_TOP = Inches(0.5)
FOTO2_WIDTH = Inches(3.5)
FOTO2_HEIGHT = Inches(3.5)

# ============================================================================
# FUNGSI HELPER
# ============================================================================

def print_header(text):
    """Print header dengan garis"""
    print("\n" + "=" * 70)
    print(text)
    print("=" * 70)

def download_from_drive(file_id, save_path, max_retries=3):
    """Download file dari Google Drive"""
    if not file_id or pd.isna(file_id):
        return False
    
    url = f"https://drive.google.com/uc?export=download&id={file_id}"
    
    for attempt in range(max_retries):
        try:
            session = requests.Session()
            response = session.get(url, timeout=30, stream=True)
            
            if response.status_code == 200:
                with open(save_path, 'wb') as f:
                    for chunk in response.iter_content(chunk_size=8192):
                        if chunk:
                            f.write(chunk)
                
                # Verifikasi ukuran file
                if os.path.getsize(save_path) > 1000:
                    return True
                else:
                    os.remove(save_path)
                    return False
            else:
                if attempt < max_retries - 1:
                    sleep(1)
                    continue
                return False
                
        except Exception as e:
            if attempt < max_retries - 1:
                sleep(1)
                continue
            print(f"   Error: {str(e)[:50]}")
            return False
    
    return False

def resize_and_crop_image(image_path, target_width, target_height):
    """Resize dan crop image menjadi persegi dengan latar merah"""
    try:
        img = Image.open(image_path)
        
        # Convert to RGB if needed
        if img.mode != 'RGB':
            img = img.convert('RGB')
        
        # Hitung rasio untuk crop ke square (4:6 photo ratio)
        width, height = img.size
        target_ratio = 4 / 6
        current_ratio = width / height
        
        if current_ratio > target_ratio:
            # Width lebih besar, crop width
            new_width = int(height * target_ratio)
            left = (width - new_width) // 2
            img = img.crop((left, 0, left + new_width, height))
        else:
            # Height lebih besar, crop height
            new_height = int(width / target_ratio)
            top = (height - new_height) // 2
            img = img.crop((0, top, width, top + new_height))
        
        # Resize untuk PowerPoint
        img.thumbnail((800, 1200), Image.Resampling.LANCZOS)
        
        # Save as temp file
        temp_path = image_path.replace('.jpg', '_processed.jpg')
        img.save(temp_path, 'JPEG', quality=90)
        
        return temp_path
    except Exception as e:
        print(f"   Error processing image: {str(e)}")
        return image_path

def add_mahasiswa_to_slide(slide, foto_path, nama, tempat_lahir, tanggal_lahir, 
                           position='left'):
    """Tambahkan data mahasiswa ke slide"""
    
    # Tentukan posisi
    if position == 'left':
        foto_left = FOTO1_LEFT
        foto_top = FOTO1_TOP
        foto_width = FOTO1_WIDTH
        foto_height = FOTO1_HEIGHT
    else:
        foto_left = FOTO2_LEFT
        foto_top = FOTO2_TOP
        foto_width = FOTO2_WIDTH
        foto_height = FOTO2_HEIGHT
    
    # Tambahkan foto
    if foto_path and os.path.exists(foto_path):
        try:
            # Process image
            processed_path = resize_and_crop_image(foto_path, foto_width, foto_height)
            
            # Add to slide
            slide.shapes.add_picture(
                processed_path,
                foto_left, foto_top,
                width=foto_width, height=foto_height
            )
            
            # Cleanup processed image
            if processed_path != foto_path and os.path.exists(processed_path):
                os.remove(processed_path)
                
        except Exception as e:
            print(f"   Error adding photo: {str(e)}")
            # Add placeholder rectangle
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                foto_left, foto_top, foto_width, foto_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(220, 220, 220)
    
    # Tambahkan nama (kotak putih)
    nama_top = foto_top + foto_height
    nama_height = Inches(0.6)
    
    nama_box = slide.shapes.add_shape(
        1,  # Rectangle
        foto_left, nama_top, foto_width, nama_height
    )
    nama_box.fill.solid()
    nama_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    nama_box.line.color.rgb = RGBColor(0, 0, 0)
    nama_box.line.width = Pt(1)
    
    # Text nama
    text_frame = nama_box.text_frame
    text_frame.clear()
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    text_frame.vertical_anchor = 1  # Middle
    
    p = text_frame.paragraphs[0]
    p.text = nama
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.name = 'Arial'
    
    # Tambahkan tempat & tanggal lahir (kotak biru)
    ttl_top = nama_top + nama_height
    ttl_height = Inches(0.8)
    
    ttl_box = slide.shapes.add_shape(
        1,  # Rectangle
        foto_left, ttl_top, foto_width, ttl_height
    )
    ttl_box.fill.solid()
    ttl_box.fill.fore_color.rgb = RGBColor(0, 102, 153)  # Biru tua
    ttl_box.line.color.rgb = RGBColor(0, 0, 0)
    ttl_box.line.width = Pt(1)
    
    # Text TTL
    text_frame = ttl_box.text_frame
    text_frame.clear()
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    text_frame.vertical_anchor = 1  # Middle
    
    # Tempat lahir
    p1 = text_frame.paragraphs[0]
    p1.text = f"Tempat Lahir: {tempat_lahir}"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(12)
    p1.font.name = 'Arial'
    p1.font.color.rgb = RGBColor(255, 255, 255)
    
    # Tanggal lahir
    p2 = text_frame.add_paragraph()
    p2.text = f"Tanggal Lahir: {tanggal_lahir}"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(12)
    p2.font.name = 'Arial'
    p2.font.color.rgb = RGBColor(255, 255, 255)

# ============================================================================
# MAIN PROGRAM
# ============================================================================

def main():
    """Main program"""
    
    print_header("GENERATOR SLIDESHOW YUDISIUM PPG 2025")
    
    # Check file Excel
    if not os.path.exists(EXCEL_FILE):
        print(f"\n❌ ERROR: File '{EXCEL_FILE}' tidak ditemukan!")
        print(f"   Pastikan file Excel ada di folder yang sama dengan script ini.")
        input("\nTekan ENTER untuk keluar...")
        sys.exit(1)
    
    # Load data
    print(f"\n📂 Membaca file Excel: {EXCEL_FILE}")
    df = pd.read_excel(EXCEL_FILE)
    print(f"   ✓ Total mahasiswa: {len(df)}")
    
    # Extract Google Drive ID
    import re
    def extract_drive_id(link):
        if pd.isna(link):
            return None
        try:
            match = re.search(r'id=([a-zA-Z0-9_-]+)', str(link))
            if match:
                return match.group(1)
            match = re.search(r'/d/([a-zA-Z0-9_-]+)', str(link))
            if match:
                return match.group(1)
        except:
            pass
        return None
    
    df['Drive_ID'] = df['File Pas Foto (Terbaru Sesuai Contoh)'].apply(extract_drive_id)
    
    # Format tanggal
    def format_tanggal(date_val):
        try:
            if pd.isna(date_val):
                return ""
            if isinstance(date_val, str):
                return date_val
            months = ['Januari', 'Februari', 'Maret', 'April', 'Mei', 'Juni',
                      'Juli', 'Agustus', 'September', 'Oktober', 'November', 'Desember']
            return f"{date_val.day:02d} {months[date_val.month - 1]} {date_val.year}"
        except:
            return str(date_val)
    
    df['Tanggal_Lahir'] = df['Tanggal Lahir'].apply(format_tanggal)
    
    # Create foto folder
    os.makedirs(FOTO_FOLDER, exist_ok=True)
    
    # STEP 1: Download foto
    print_header("STEP 1: DOWNLOAD FOTO DARI GOOGLE DRIVE")
    print(f"📁 Folder tujuan: {FOTO_FOLDER}")
    print(f"📥 Total foto yang akan di-download: {len(df)}")
    
    response = input(f"\n⚠️  Download akan memakan waktu ~20-30 menit.\n    Lanjutkan? (y/n): ")
    if response.lower() != 'y':
        print("\n❌ Proses dibatalkan.")
        input("Tekan ENTER untuk keluar...")
        sys.exit(0)
    
    print(f"\n⏳ Proses download dimulai...")
    print("-" * 70)
    
    success = 0
    failed = 0
    
    for idx, row in df.iterrows():
        no = row['No']
        nama = row['Nama Lengkap Sesuai Ijazah']
        drive_id = row['Drive_ID']
        
        filename = f"{no:04d}.jpg"
        filepath = os.path.join(FOTO_FOLDER, filename)
        
        # Skip if already downloaded
        if os.path.exists(filepath):
            success += 1
            continue
        
        # Progress
        if (idx + 1) % 50 == 0:
            print(f"Progress: {idx + 1}/{len(df)} ({(idx+1)/len(df)*100:.1f}%) - "
                  f"Success: {success}, Failed: {failed}")
        
        # Download
        if download_from_drive(drive_id, filepath):
            success += 1
        else:
            failed += 1
        
        # Small delay
        sleep(0.3)
    
    print("-" * 70)
    print(f"✓ Download selesai!")
    print(f"  Berhasil: {success}/{len(df)}")
    print(f"  Gagal: {failed}/{len(df)}")
    
    if failed > 0:
        print(f"\n⚠️  {failed} foto gagal di-download. PowerPoint tetap akan dibuat,")
        print(f"    tapi slide untuk foto yang gagal akan kosong.")
    
    # STEP 2: Generate PowerPoint
    print_header("STEP 2: GENERATE POWERPOINT")
    print(f"📊 Total slide yang akan dibuat: {(len(df) + 1) // 2}")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Add blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout
    
    print(f"\n⏳ Membuat slide...")
    
    # Generate slides (2 mahasiswa per slide)
    for i in range(0, len(df), 2):
        # Add new slide
        slide = prs.slides.add_slide(blank_layout)
        
        # Set background putih
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Mahasiswa 1 (kiri)
        row1 = df.iloc[i]
        foto1_path = os.path.join(FOTO_FOLDER, f"{row1['No']:04d}.jpg")
        add_mahasiswa_to_slide(
            slide, foto1_path,
            row1['Nama Lengkap Sesuai Ijazah'],
            row1['Tempat Lahir'],
            row1['Tanggal_Lahir'],
            position='left'
        )
        
        # Mahasiswa 2 (kanan) - jika ada
        if i + 1 < len(df):
            row2 = df.iloc[i + 1]
            foto2_path = os.path.join(FOTO_FOLDER, f"{row2['No']:04d}.jpg")
            add_mahasiswa_to_slide(
                slide, foto2_path,
                row2['Nama Lengkap Sesuai Ijazah'],
                row2['Tempat Lahir'],
                row2['Tanggal_Lahir'],
                position='right'
            )
        
        # Progress
        if (i + 2) % 100 == 0:
            slide_no = (i // 2) + 1
            total_slides = (len(df) + 1) // 2
            print(f"Progress: {slide_no}/{total_slides} slide ({slide_no/total_slides*100:.1f}%)")
    
    # Save PowerPoint
    print(f"\n💾 Menyimpan PowerPoint: {OUTPUT_PPTX}")
    prs.save(OUTPUT_PPTX)
    
    print_header("✅ SELESAI!")
    print(f"\n📄 PowerPoint berhasil dibuat: {OUTPUT_PPTX}")
    print(f"📊 Total slide: {len(prs.slides)}")
    print(f"📁 Foto tersimpan di folder: {FOTO_FOLDER}")
    print(f"\n🎉 Silakan buka file PowerPoint dan cek hasilnya!")
    
    input("\n\nTekan ENTER untuk keluar...")

if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\n\n❌ Proses dibatalkan oleh user.")
        input("Tekan ENTER untuk keluar...")
    except Exception as e:
        print(f"\n\n❌ ERROR: {str(e)}")
        import traceback
        traceback.print_exc()
        input("\nTekan ENTER untuk keluar...")
